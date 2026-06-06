from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

TEMPLATE_PATH = "data/templates/default.pptx"

def create_starter_template():
    """Create a stylish starter template if it doesn't exist."""
    os.makedirs("data/templates", exist_ok=True)
    
    prs = Presentation()
    
    # Customize Master Slide (basic simulation by editing layouts)
    # We will just edit the first few meaningful layouts to have a background
    
    # 1. Title Slide Layout
    slide_layout = prs.slide_master.slide_layouts[0]
    background = slide_layout.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102) # Dark Blue background
    
    # Customize Title Text Style
    title_style = slide_layout.placeholders[0].text_frame.paragraphs[0].font
    title_style.name = "Arial"
    title_style.size = Pt(44)
    title_style.color.rgb = RGBColor(255, 255, 255) # White text
    
    # Customize Subtitle
    subtitle = slide_layout.placeholders[1].text_frame.paragraphs[0].font
    subtitle.name = "Arial"
    subtitle.color.rgb = RGBColor(200, 200, 200) # Light gray
    
    # 2. Content Slide Layout
    bullet_layout = prs.slide_master.slide_layouts[1]
    background = bullet_layout.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255) # AliceBlue (Very light)
    
    # Title on content slide
    title_c = bullet_layout.placeholders[0].text_frame.paragraphs[0].font
    title_c.name = "Arial"
    title_c.color.rgb = RGBColor(0, 51, 102) # Dark Blue
    title_c.bold = True
    
    # Body text
    # We can't easily set default font for all body placeholders in python-pptx without iterating,
    # but setting the master layout helps new slides inherit.
    
    prs.save(TEMPLATE_PATH)
    print(f"Starter template created at {TEMPLATE_PATH}")

if __name__ == "__main__":
    create_starter_template()
