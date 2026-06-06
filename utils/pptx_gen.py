from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches, Pt
import os
from utils.template_gen import create_starter_template, TEMPLATE_PATH

def create_presentation_pptx(topic: str, content: str, output_path: str, template_path: str = None):
    if not template_path:
        if not os.path.exists(TEMPLATE_PATH):
            create_starter_template()
        template_path = TEMPLATE_PATH
    
    try:
        prs = Presentation(template_path)
    except Exception:
        # Fallback if template is corrupt or not found
        prs = Presentation()
    
    # Title Slide
    try:
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = topic
            
        try:
            subtitle = slide.placeholders[1]
            subtitle.text = "Talaba Bot tomonidan tayyorlandi"
        except Exception:
            pass
    except Exception:
        pass
    
    import re
    import logging
    import json
    logger = logging.getLogger(__name__)

    parsed_slides = []
    
    # Try to parse JSON first (new format)
    try:
        # Find JSON array in the string
        json_start = content.find('[')
        json_end = content.rfind(']') + 1
        if json_start != -1 and json_end != -1:
            json_str = content[json_start:json_end]
            parsed_slides = json.loads(json_str)
    except Exception as e:
        logger.warning(f"PPTX Gen: JSON parsing failed: {e}")

    # Fallback to old format if JSON failed
    if not parsed_slides:
        sections = content.split("|||")
        if len(sections) < 2:
            sections = re.split(r'(?:^|\n)(?:Slayd|Slide)\s*\d+[:\.]?', content, flags=re.IGNORECASE)
        if len(sections) < 2:
            sections = content.split("\n\n")

        for section in sections:
            section = section.strip()
            if len(section) < 5: continue
            
            lines = section.split("\n", 1)
            slide_title = lines[0].strip()
            slide_text = lines[1].strip() if len(lines) > 1 else ""
            
            if not slide_text and len(slide_title) > 50:
                 parts = re.split(r'[:.?!]\s', slide_title, 1)
                 if len(parts) > 1:
                     slide_title = parts[0]
                     slide_text = parts[1]
                     
            parsed_slides.append({"title": slide_title, "content": slide_text})

    valid_slides_count = 0
    for slide_data in parsed_slides:
        slide_title = slide_data.get("title", "")
        slide_text = slide_data.get("content", "")
        
        # Cleanup
        slide_title = re.sub(r'^(?:Slayd|Slide|Step|Bo\'lim)\s*\d*[:\.]?\s*', '', slide_title, flags=re.IGNORECASE).strip()
        slide_title = re.sub(r'^\d+[:\.]\s*', '', slide_title).strip()
        slide_title = slide_title.replace('**', '').replace('__', '').replace('*', '')
        
        # Remove markdown from text too
        slide_text = slide_text.replace('**', '').replace('__', '').replace('*', '')

        try:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            if shapes.title:
                shapes.title.text = slide_title
                
            try:
                body_shape = shapes.placeholders[1]
                body_shape.text_frame.text = slide_text
                # Intentionally not setting fixed font size here to allow PowerPoint Auto-fit to work!
            except Exception:
                # If placeholder 1 doesn't exist, we just skip text insertion
                pass
                
            valid_slides_count += 1
        except Exception as e:
            logger.error(f"Error adding slide: {e}")

    if valid_slides_count == 0:
        # Emergency fallback
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.shapes.title: slide.shapes.title.text = "Taqdimot"
            slide.shapes.placeholders[1].text_frame.text = content[:1000]
        except:
            pass

    prs.save(output_path)
    return output_path
